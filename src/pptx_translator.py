from typing import Optional
from pptx import Presentation
from dotenv import load_dotenv
import os
import anthropic
import re
from pathlib import Path
from functools import lru_cache
from tqdm import tqdm
from controler import language_codes

translation_cache = {}

parent_dir = Path(__file__).parent.parent
load_dotenv(dotenv_path=parent_dir / '.env')
API_KEY_ANTHROPIC = os.getenv('API_KEY_ANTHROPIC')
client = anthropic.Anthropic(
    api_key=API_KEY_ANTHROPIC,
)

def is_exception_text(text: str, source_lang: str, target_lang: str, version: str) -> Optional[str]:
    """
    Check if the given text contains exception patterns and modify it accordingly.

    This function checks for two types of exceptions:
    1. Language code: Replaces the uppercase source language code with the uppercase target language code.
    2. Version number: Replaces 'V.' followed by three digits with 'v' followed by the digits from the version parameter.

    Args:
        text (str): The input text to check for exceptions.
        source_lang (str): The source language code (e.g., 'en' for English).
        target_lang (str): The target language code (e.g., 'fr' for French).
        version (str): The version string to use for replacement.

    Returns:
        Optional[str]: 
            - If an exception is found and modified, returns the modified text.
            - If no exception is found, returns None.
    """
    # Format version to 'v' followed by digits from the version
    version_formatted = 'v' + ''.join(c for c in version if c.isdigit())
    
    exceptions = [
        f"- {source_lang.upper()}",
        r"V\.\d{3}"  # Regex to match V. followed by 3 digits
    ]
    
    for exception in exceptions:
        if exception == f"- {source_lang.upper()}" and exception in text:
            return text.replace(f"- {source_lang.upper()}", f"- {target_lang.upper()}")
        elif re.search(exception, text):  # Use regex search for version
            return re.sub(r"V\.\d{3}", version_formatted, text)
    
    return None

is_exception_text()
def translate_txt_to(text, language):
    try:
        message = client.messages.create(
            model="claude-3-5-sonnet-20240620",
            max_tokens=2000,
            temperature=0.2,
            system=f"Your task is to accurately translate this text into {language}. Only output the translation. If there's nothing to translate simply output the original text.",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"string to translate:\n {text}"
                        }
                    ]
                }
            ]
        )
        
        translated_text = message.content[0].text
        return translated_text

    except Exception as e:
        raise TranslationError(f"Translation failed: {str(e)}")

class TranslationError(Exception):
    pass

@lru_cache(maxsize=1000)
def get_translation(text, target_language):
    """
    Check if the translation exists in cache, if not, translate and cache the result.
    """
    cache_key = (text, target_language)
    if cache_key in translation_cache:
        return translation_cache[cache_key]
    else:
        translated_text = translate_txt_to(text, target_language)
        translation_cache[cache_key] = translated_text
        return translated_text

def count_total_runs(prs):
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    total += len(paragraph.runs)
    return total

def translate_pptx(input_path, output_path, source_lang, target_lang, version, use_exception=False):
    prs = Presentation(input_path)
    total_runs = count_total_runs(prs)
    progress_bar = tqdm(total=total_runs, unit='run', desc='Translating')

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if use_exception:
                        exception_result = is_exception_text(run.text, source_lang, target_lang, version)
                        if exception_result:
                            run.text = exception_result
                            progress_bar.update(1)
                            continue

                    translated_text = get_translation(run.text, language_codes[target_lang])
                    run.text = translated_text
                    progress_bar.update(1)

    progress_bar.close()
    prs.save(output_path)
